# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Primary color: #17375E = RGB(23, 55, 94)
PRIMARY_COLOR = (23, 55, 94)
ACCENT_COLOR = (46, 134, 193)  # Lighter blue for contrast
SUCCESS_COLOR = (39, 174, 96)
WARNING_COLOR = (230, 126, 34)
TEXT_COLOR = (50, 50, 50)
WHITE = (255, 255, 255)
LIGHT_BG = (245, 247, 250)


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_title_slide(prs, title_en, title_vi, subtitle_en, subtitle_vi):
    """Title slide with bilingual text"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(*PRIMARY_COLOR)
    bg.line.fill.background()

    # Title English
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = title_en
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*WHITE)
    p.alignment = PP_ALIGN.CENTER

    # Title Vietnamese
    box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = title_vi
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(200, 210, 220)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle English
    box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = subtitle_en
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(180, 190, 200)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle Vietnamese
    box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12.333), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = subtitle_vi
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(160, 170, 180)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title_en, title_vi, items):
    """Content slide with bilingual title and items (each item is [en, vi])"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title English
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = title_en
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*PRIMARY_COLOR)

    # Title Vietnamese
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(12.333), Inches(0.5))
    p = box.text_frame.paragraphs[0]
    p.text = title_vi
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(*ACCENT_COLOR)

    # Underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.35), Inches(2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(*PRIMARY_COLOR)
    line.line.fill.background()

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if isinstance(item, list) and len(item) == 2:
            # Bilingual item [en, vi]
            en_text, vi_text = item

            # English line
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = en_text
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*TEXT_COLOR)
            p.space_after = Pt(2)

            # Vietnamese line
            p = tf.add_paragraph()
            p.text = vi_text
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(100, 100, 100)
            p.space_after = Pt(14)
        else:
            # Single line (spacer or special)
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(item) if item else ""
            p.font.size = Pt(18)
            p.space_after = Pt(8)

    return slide


def add_table_slide(prs, title_en, title_vi, headers, rows):
    """Table slide with bilingual title"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title English
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = title_en
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*PRIMARY_COLOR)

    # Title Vietnamese
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(12.333), Inches(0.5))
    p = box.text_frame.paragraphs[0]
    p.text = title_vi
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(*ACCENT_COLOR)

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_width = Inches(12.333)

    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.5), table_width, Inches(0.45 * num_rows)).table

    # Headers
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*PRIMARY_COLOR)
        for para in cell.text_frame.paragraphs:
            para.font.bold = True
            para.font.color.rgb = RGBColor(*WHITE)
            para.font.size = Pt(14)

    # Rows
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(cell_text)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*LIGHT_BG)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(12)
                para.font.color.rgb = RGBColor(*TEXT_COLOR)

    return slide


def add_two_column_slide(prs, title_en, title_vi, left_title, left_items, right_title, right_items):
    """Two column comparison slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title English
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = title_en
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*PRIMARY_COLOR)

    # Title Vietnamese
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(12.333), Inches(0.5))
    p = box.text_frame.paragraphs[0]
    p.text = title_vi
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(*ACCENT_COLOR)

    # Left column title
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.5))
    p = box.text_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*SUCCESS_COLOR)

    # Left column content
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.95), Inches(5.8), Inches(5))
    tf = box.text_frame
    for i, item in enumerate(left_items):
        if isinstance(item, list) and len(item) == 2:
            # English
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item[0]
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*TEXT_COLOR)
            p.space_after = Pt(1)
            # Vietnamese
            p = tf.add_paragraph()
            p.text = item[1]
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(100, 100, 100)
            p.space_after = Pt(8)
        else:
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(item)
            p.font.size = Pt(15)
            p.space_after = Pt(6)

    # Right column title
    box = slide.shapes.add_textbox(Inches(7), Inches(1.4), Inches(5.8), Inches(0.5))
    p = box.text_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*WARNING_COLOR)

    # Right column content
    box = slide.shapes.add_textbox(Inches(7), Inches(1.95), Inches(5.8), Inches(5))
    tf = box.text_frame
    for i, item in enumerate(right_items):
        if isinstance(item, list) and len(item) == 2:
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item[0]
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*TEXT_COLOR)
            p.space_after = Pt(1)
            p = tf.add_paragraph()
            p.text = item[1]
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(100, 100, 100)
            p.space_after = Pt(8)
        else:
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(item)
            p.font.size = Pt(15)
            p.space_after = Pt(6)

    return slide


def add_code_slide(prs, title_en, title_vi, code_text):
    """Code slide with bilingual title"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title English
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.5))
    p = box.text_frame.paragraphs[0]
    p.text = title_en
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*PRIMARY_COLOR)

    # Title Vietnamese
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(12.333), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = title_vi
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(*ACCENT_COLOR)

    # Code background
    code_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.333), Inches(5.8))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = RGBColor(40, 44, 52)
    code_bg.line.fill.background()

    # Code text
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.5))
    p = box.text_frame.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(14)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(220, 220, 220)

    return slide


def add_qa_slide(prs, resources=None):
    """Q&A slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(*PRIMARY_COLOR)
    bg.line.fill.background()

    # Q&A
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = "Questions & Answers"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*WHITE)
    p.alignment = PP_ALIGN.CENTER

    box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.333), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = "Hỏi Đáp"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(180, 190, 200)
    p.alignment = PP_ALIGN.CENTER

    if resources:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(2))
        tf = box.text_frame
        for i, res in enumerate(resources):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = res
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(160, 170, 180)
            p.alignment = PP_ALIGN.CENTER

    return slide


# ============================================
# PRESENTATION 1: ASTRO + SANITY
# ============================================
def create_astro_sanity_pptx():
    prs = create_presentation()

    # Slide 1: Title
    add_title_slide(prs,
        "Astro + Sanity",
        "Kết hợp Astro và Sanity",
        "Modern Content-Driven Architecture",
        "Kiến trúc hiện đại hướng nội dung"
    )

    # Slide 2: Agenda
    add_content_slide(prs, "Agenda", "Nội dung trình bày", [
        ["1. Problems with traditional CMS", "Vấn đề với CMS truyền thống"],
        ["2. Introduction to Sanity", "Giới thiệu về Sanity"],
        ["3. Introduction to Astro", "Giới thiệu về Astro"],
        ["4. Why combine Sanity + Astro?", "Tại sao kết hợp Sanity + Astro?"],
        ["5. Comparisons & Use cases", "So sánh và trường hợp sử dụng"]
    ])

    # Slide 3: Problems
    add_content_slide(prs, "Problems with Traditional CMS", "Vấn đề với CMS truyền thống", [
        ["Poor performance (WordPress, Drupal)", "Hiệu suất kém (WordPress, Drupal)"],
        ["Difficult to customize", "Khó tùy chỉnh theo nhu cầu"],
        ["Monolithic - Frontend + Backend tightly coupled", "Monolithic - Frontend + Backend gắn chặt"],
        ["Hard to scale when traffic increases", "Khó mở rộng khi traffic tăng"],
        ["Vendor lock-in", "Bị ràng buộc vào một nền tảng"]
    ])

    # Slide 4: What is Headless CMS
    add_content_slide(prs, "What is Headless CMS?", "Headless CMS là gì?", [
        ["Separate Backend (CMS) and Frontend", "Tách biệt Backend (CMS) và Frontend"],
        ["Content delivered via API (REST/GraphQL)", "Nội dung cung cấp qua API"],
        ["Frontend can be any framework", "Frontend có thể là bất kỳ framework nào"],
        ["Content = Data, presentation independent", "Content = Data, không phụ thuộc hiển thị"],
        ["Multi-platform: Web, Mobile, IoT", "Đa nền tảng: Web, Mobile, IoT"]
    ])

    # Slide 5: What is Sanity
    add_content_slide(prs, "What is Sanity?", "Sanity là gì?", [
        ["Headless CMS - Structured Content Platform", "Headless CMS - Nền tảng nội dung có cấu trúc"],
        ["Real-time collaboration (like Google Docs)", "Cộng tác thời gian thực (như Google Docs)"],
        ["GROQ - Powerful query language", "GROQ - Ngôn ngữ truy vấn mạnh mẽ"],
        ["Sanity Studio - Fully customizable admin UI", "Sanity Studio - Giao diện admin tùy chỉnh hoàn toàn"],
        ["Content Lake - Global CDN", "Content Lake - CDN toàn cầu"]
    ])

    # Slide 6: Sanity Strengths/Weaknesses
    add_two_column_slide(prs,
        "Sanity - Strengths vs Weaknesses",
        "Sanity - Điểm mạnh vs Điểm yếu",
        "Strengths / Điểm mạnh",
        [
            ["Generous free tier (20GB, 10K docs)", "Free tier hào phóng"],
            ["Real-time collaboration", "Cộng tác thời gian thực"],
            ["Flexible code-based schema", "Schema linh hoạt (code-based)"],
            ["Live preview & visual editing", "Xem trước trực tiếp"],
            ["Global CDN - great performance", "CDN toàn cầu - hiệu suất tốt"]
        ],
        "Weaknesses / Điểm yếu",
        [
            ["Learning curve (GROQ, schema)", "Cần thời gian học (GROQ, schema)"],
            ["Pricing increases at scale", "Giá tăng khi mở rộng"],
            ["Requires developer to setup", "Cần developer để cài đặt"],
            ["No built-in frontend", "Không có frontend sẵn"],
            ["Vendor lock-in (GROQ)", "Ràng buộc vendor (GROQ)"]
        ]
    )

    # Slide 7: CMS Comparison
    add_table_slide(prs,
        "Headless CMS Comparison",
        "So sánh các Headless CMS",
        ["Feature", "Sanity", "Contentful", "Strapi", "WP Headless"],
        [
            ["Pricing", "Free 20GB", "Free 5 users", "Free self-host", "Free"],
            ["Schema", "Code-based", "UI-based", "UI + Code", "Predefined"],
            ["Query", "GROQ", "GraphQL", "REST/GraphQL", "REST/GraphQL"],
            ["Real-time", "Yes", "No (Webhook)", "No", "No"],
            ["Self-host", "No", "No", "Yes", "Yes"],
            ["Customization", "High", "Medium", "High", "Medium"]
        ]
    )

    # Slide 8: What is Astro
    add_content_slide(prs, "What is Astro?", "Astro là gì?", [
        ["Static Site Generator + SSR hybrid", "Tạo trang tĩnh + SSR kết hợp"],
        ["Island Architecture - Partial hydration", "Kiến trúc Island - Hydration một phần"],
        ["Zero JavaScript by default", "Không JavaScript mặc định"],
        ["Multi-framework (React, Vue, Svelte)", "Hỗ trợ nhiều framework"],
        ["Content Collections built-in", "Content Collections tích hợp sẵn"],
        ["Optimized for content-driven websites", "Tối ưu cho website hướng nội dung"]
    ])

    # Slide 9: Astro Strengths/Weaknesses
    add_two_column_slide(prs,
        "Astro - Strengths vs Weaknesses",
        "Astro - Điểm mạnh vs Điểm yếu",
        "Strengths / Điểm mạnh",
        [
            ["Extremely fast (Zero JS)", "Cực nhanh (Zero JS)"],
            ["Islands - hydrate only when needed", "Islands - chỉ hydrate khi cần"],
            ["Multi-framework in one project", "Nhiều framework trong 1 project"],
            ["Easy to learn", "Dễ học"],
            ["Fast build speed", "Tốc độ build nhanh"]
        ],
        "Weaknesses / Điểm yếu",
        [
            ["Smaller ecosystem than Next.js", "Ecosystem nhỏ hơn Next.js"],
            ["Fewer tutorials/resources", "Ít tài liệu hướng dẫn"],
            ["SSR still maturing", "SSR còn mới"],
            ["Not suitable for complex apps", "Không phù hợp app phức tạp"],
            ["Less enterprise adoption", "Ít doanh nghiệp sử dụng"]
        ]
    )

    # Slide 10: Framework Comparison
    add_table_slide(prs,
        "Framework Comparison",
        "So sánh các Framework",
        ["Feature", "Astro", "Next.js", "Nuxt", "Gatsby"],
        [
            ["JS shipped", "Minimal", "Full bundle", "Full bundle", "Full bundle"],
            ["Build speed", "Fast", "Medium", "Medium", "Slow"],
            ["SSG + SSR", "Yes", "Yes", "Yes", "Limited"],
            ["Multi-framework", "Yes", "React only", "Vue only", "React only"],
            ["Learning curve", "Easy", "Medium", "Medium", "Hard"],
            ["Best for", "Content sites", "Full apps", "Vue apps", "Data-heavy"]
        ]
    )

    # Slide 11: Why Sanity + Astro
    add_content_slide(prs, "Why Sanity + Astro?", "Tại sao Sanity + Astro?", [
        ["Sanity: Structured content, real-time editing", "Sanity: Nội dung có cấu trúc, chỉnh sửa thời gian thực"],
        ["Astro: Fast static pages, zero JS default", "Astro: Trang tĩnh nhanh, không JS mặc định"],
        ["Sanity: Live preview integration", "Sanity: Tích hợp xem trước trực tiếp"],
        ["Astro: Island hydration for interactivity", "Astro: Island hydration cho tương tác"],
        ["", ""],
        ["= Content-first, Performance-optimized", "= Ưu tiên nội dung, Tối ưu hiệu suất"]
    ])

    # Slide 12: Architecture
    add_code_slide(prs, "Architecture", "Kiến trúc hệ thống", """
+------------------+      +------------------+      +------------------+
|                  |      |                  |      |                  |
|  Sanity Studio   |----->|   Astro Build    |----->|  Vercel / CDN    |
|  (CMS)           |      |   (Frontend)     |      |  (Deploy)        |
|                  |      |                  |      |                  |
+------------------+      +------------------+      +------------------+
        |                         |
        v                         v
   Content Lake              Static HTML
   (Real-time)               + Islands


Workflow:
1. Content Editor --> Sanity Studio
2. Developer --> Astro frontend
3. Preview --> Live preview integration
4. Deploy --> Vercel/Netlify auto-build
5. Update --> Webhook trigger rebuild
""")

    # Slide 13: Use Cases
    add_content_slide(prs, "Suitable Use Cases", "Trường hợp sử dụng phù hợp", [
        ["Marketing websites", "Website marketing"],
        ["Documentation sites", "Website tài liệu"],
        ["Blogs & magazines", "Blog và tạp chí"],
        ["Landing pages", "Trang landing"],
        ["Multi-language sites", "Website đa ngôn ngữ"],
        ["E-commerce (content-heavy)", "Thương mại điện tử (nhiều nội dung)"]
    ])

    # Slide 14: Q&A
    add_qa_slide(prs, [
        "sanity.io",
        "astro.build",
        "docs.astro.build/guides/cms/sanity"
    ])

    prs.save("01-Astro-Sanity.pptx")
    print("[OK] Created: 01-Astro-Sanity.pptx")


# ============================================
# PRESENTATION 2: CLAUDE PLUGINS
# ============================================
def create_claude_plugins_pptx():
    prs = create_presentation()

    # Slide 1: Title
    add_title_slide(prs,
        "Claude Plugins",
        "Plugin cho Claude",
        "Extend Claude Code with Skills, Agents & MCP",
        "Mở rộng Claude Code với Skills, Agents và MCP"
    )

    # Slide 2: Agenda
    add_content_slide(prs, "Agenda", "Nội dung trình bày", [
        ["1. What are Claude Plugins?", "Claude Plugins là gì?"],
        ["2. Plugin Components", "Các thành phần của Plugin"],
        ["3. Plugin Structure", "Cấu trúc Plugin"],
        ["4. How to Build", "Cách xây dựng"],
        ["5. Strengths & Weaknesses", "Điểm mạnh và điểm yếu"],
        ["6. Best Practices", "Các thực hành tốt nhất"]
    ])

    # Slide 3: What is Claude Plugins
    add_content_slide(prs, "What are Claude Plugins?", "Claude Plugins là gì?", [
        ["Bundle tools, skills, integrations for one-click install", "Gói công cụ, kỹ năng, tích hợp cài đặt 1-click"],
        ["Extensions for Claude Code & Claude Cowork", "Phần mở rộng cho Claude Code & Cowork"],
        ["10,000+ plugins available in marketplace", "10,000+ plugin có sẵn trên marketplace"],
        ["Extend AI capabilities", "Mở rộng khả năng AI"]
    ])

    # Slide 4: Plugin Components
    add_table_slide(prs,
        "Plugin Components",
        "Các thành phần của Plugin",
        ["Component", "Description EN", "Mô tả VI"],
        [
            ["Skills", "Auto-apply instructions", "Hướng dẫn tự động áp dụng"],
            ["Commands", "Slash commands /cmd", "Lệnh slash /cmd"],
            ["Agents", "Custom sub-agents", "Sub-agents tùy chỉnh"],
            ["Hooks", "Event handlers", "Xử lý sự kiện"],
            ["MCP", "External API integration", "Tích hợp API bên ngoài"],
            ["LSP", "Code intelligence", "Thông minh mã nguồn"]
        ]
    )

    # Slide 5: Skills vs MCP
    add_two_column_slide(prs,
        "Skills vs MCP Servers",
        "So sánh Skills và MCP Servers",
        "Skills",
        [
            ["Knowledge, guidelines, workflows", "Kiến thức, hướng dẫn, quy trình"],
            ["Claude auto-detects context", "Claude tự nhận biết context"],
            ["Static instructions", "Hướng dẫn tĩnh"],
            ["Example: Code review rules", "Ví dụ: Quy tắc review code"]
        ],
        "MCP Servers",
        [
            ["External data, APIs", "Dữ liệu bên ngoài, APIs"],
            ["Explicit tool calls", "Gọi tool cụ thể"],
            ["Dynamic real-time data", "Dữ liệu động thời gian thực"],
            ["Example: Query database", "Ví dụ: Truy vấn database"]
        ]
    )

    # Slide 6: Directory Structure
    add_code_slide(prs, "Plugin Directory Structure", "Cấu trúc thư mục Plugin", """
my-plugin/
|
+-- .claude-plugin/
|   +-- plugin.json        # Manifest (REQUIRED / BẮT BUỘC)
|
+-- skills/                # Model-invoked skills
|   +-- code-review/
|       +-- SKILL.md
|
+-- commands/              # Slash commands
|   +-- deploy.md
|
+-- agents/                # Custom agents
|   +-- reviewer.md
|
+-- hooks/
|   +-- hooks.json         # Event handlers
|
+-- .mcp.json              # MCP server configs
|
+-- README.md


NOTE: Directories at ROOT level, NOT inside .claude-plugin/
LƯU Ý: Các thư mục ở ROOT, KHÔNG đặt trong .claude-plugin/
""")

    # Slide 7: Plugin Manifest
    add_code_slide(prs, "Plugin Manifest", "File plugin.json", """
{
  "name": "my-plugin",
  "description": "Short description of the plugin",
  "version": "1.0.0",
  "author": {
    "name": "Your Name",
    "email": "email@example.com"
  },
  "homepage": "https://...",
  "repository": "https://github.com/...",
  "license": "MIT"
}


KEY POINTS / ĐIỂM CHÍNH:

- name --> Namespace for skills: /my-plugin:hello
         Namespace cho skills: /my-plugin:hello

- version --> Semantic versioning (major.minor.patch)
              Phiên bản ngữ nghĩa

- description --> Shown in plugin manager
                  Hiển thị trong plugin manager
""")

    # Slide 8: Writing Skills
    add_code_slide(prs, "Writing Skills (SKILL.md)", "Viết Skills (SKILL.md)", """
---
name: code-review
description: Review code for security issues and best practices.
             Use when checking PRs or reviewing code changes.
---

When reviewing code, check for:

1. SQL injection vulnerabilities
2. XSS (Cross-site scripting)
3. Authentication issues
4. Input validation
5. Error handling


IMPORTANT / QUAN TRỌNG:

- description = trigger mechanism!
  description = cơ chế kích hoạt!

- Write clear, specific descriptions
  Viết mô tả rõ ràng, cụ thể

- Claude will auto-apply when context matches
  Claude tự động áp dụng khi context phù hợp
""")

    # Slide 9: Writing Hooks
    add_code_slide(prs, "Writing Hooks", "Viết Hooks", """
// hooks/hooks.json

{
  "hooks": {
    "PostToolUse": [
      {
        "matcher": "Write|Edit",
        "hooks": [
          {
            "type": "command",
            "command": "npm run lint:fix"
          }
        ]
      }
    ]
  }
}


AVAILABLE EVENTS / CÁC SỰ KIỆN:

- PreToolUse   : Before tool runs / Trước khi tool chạy
- PostToolUse  : After tool runs / Sau khi tool chạy
- Notification : On notification / Khi có thông báo
- Stop         : Session ends / Khi kết thúc session
""")

    # Slide 10: Installation
    add_code_slide(prs, "Installation & Testing", "Cài đặt và Kiểm tra", """
# From official marketplace / Từ marketplace chính thức
claude plugin add @anthropic/github

# From GitHub repository / Từ GitHub
claude plugin add github:username/repo-name

# From local directory / Từ thư mục local
claude plugin add ./my-plugin

# Test during development / Test trong quá trình phát triển
claude --plugin-dir ./my-plugin

# Reload after changes / Reload sau khi thay đổi
/reload-plugins


TIP: Use --plugin-dir to test before publishing
GỢI Ý: Dùng --plugin-dir để test trước khi publish
""")

    # Slide 11: Strengths/Weaknesses
    add_two_column_slide(prs,
        "Strengths vs Weaknesses",
        "Điểm mạnh vs Điểm yếu",
        "Strengths / Điểm mạnh",
        [
            ["One-click install", "Cài đặt 1-click"],
            ["Namespaced - avoid conflicts", "Namespace - tránh xung đột"],
            ["Versioned - easy updates", "Có phiên bản - dễ cập nhật"],
            ["Hot reload (/reload-plugins)", "Reload nóng"],
            ["Shareable - team/community", "Chia sẻ được"],
            ["Large ecosystem (10K+ plugins)", "Ecosystem lớn"]
        ],
        "Weaknesses / Điểm yếu",
        [
            ["Learning curve", "Cần thời gian học"],
            ["Long namespace (/plugin:skill)", "Namespace dài"],
            ["Debugging difficult", "Debug khó"],
            ["Dependency management", "Quản lý dependency"],
            ["Limited documentation", "Tài liệu còn hạn chế"],
            ["Ecosystem still new", "Ecosystem còn mới"]
        ]
    )

    # Slide 12: Standalone vs Plugin
    add_table_slide(prs,
        "Standalone vs Plugin",
        "Độc lập vs Plugin",
        ["Approach", "Skill Name", "Best For EN", "Phù hợp VI"],
        [
            ["Standalone (.claude/)", "/hello", "Personal, experiments", "Cá nhân, thử nghiệm"],
            ["Plugin", "/plugin:hello", "Share, versioned, team", "Chia sẻ, phiên bản, team"]
        ]
    )

    # Slide 13: Best Practices
    add_two_column_slide(prs,
        "Best Practices",
        "Các thực hành tốt nhất",
        "DO / NÊN",
        [
            ["Write clear, specific descriptions", "Viết mô tả rõ ràng, cụ thể"],
            ["Test with --plugin-dir first", "Test với --plugin-dir trước"],
            ["Use semantic versioning", "Dùng semantic versioning"],
            ["Include README.md", "Bao gồm README.md"],
            ["Handle errors gracefully", "Xử lý lỗi tốt"]
        ],
        "DON'T / KHÔNG NÊN",
        [
            ["Put dirs inside .claude-plugin/", "Đặt thư mục trong .claude-plugin/"],
            ["Write vague descriptions", "Viết mô tả mơ hồ"],
            ["Hardcode secrets", "Hardcode secrets"],
            ["Skip testing", "Bỏ qua testing"],
            ["Make skills too broad", "Tạo skills quá rộng"]
        ]
    )

    # Slide 14: Submit to Marketplace
    add_content_slide(prs, "Submit to Marketplace", "Gửi lên Marketplace", [
        ["Submission URLs:", "Đường dẫn gửi:"],
        ["  - claude.ai/settings/plugins/submit", ""],
        ["  - platform.claude.com/plugins/submit", ""],
        ["", ""],
        ["Requirements:", "Yêu cầu:"],
        ["  - Complete documentation", "  - Tài liệu đầy đủ"],
        ["  - Clear use case description", "  - Mô tả use case rõ ràng"],
        ["  - No security issues", "  - Không có vấn đề bảo mật"],
        ["  - Follow plugin guidelines", "  - Tuân thủ hướng dẫn plugin"]
    ])

    # Slide 15: Q&A
    add_qa_slide(prs, [
        "claude.com/plugins",
        "code.claude.com/docs/en/plugins",
        "github.com/anthropics/claude-code"
    ])

    prs.save("02-Claude-Plugins.pptx")
    print("[OK] Created: 02-Claude-Plugins.pptx")


# ============================================
# MAIN
# ============================================
if __name__ == "__main__":
    print("Creating PowerPoint presentations...\n")
    print("Primary color: #17375E\n")

    create_astro_sanity_pptx()
    create_claude_plugins_pptx()

    print("\n[DONE] Created 2 PowerPoint files:")
    print("   - 01-Astro-Sanity.pptx (14 slides)")
    print("   - 02-Claude-Plugins.pptx (15 slides)")
